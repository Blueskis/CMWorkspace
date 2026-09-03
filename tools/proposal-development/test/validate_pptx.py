#!/usr/bin/env python3
"""Independent validator for .pptx files this tool generates. Used by the
Node test suite (via subprocess) and can be run standalone:

    python3 test/validate_pptx.py out/some.pptx

Prints a small JSON summary (slide count, layout names, titles, run count
per shape) to stdout. Exits non-zero if the file will not open.
"""
import json
import sys

from pptx import Presentation


def summarize(path):
    prs = Presentation(path)
    slides = []
    for slide in prs.slides:
        shapes = []
        for shape in slide.shapes:
            entry = {"name": shape.name, "has_text_frame": shape.has_text_frame}
            if shape.has_text_frame:
                entry["paragraphs"] = [
                    {"level": p.level, "runs": [r.text for r in p.runs]}
                    for p in shape.text_frame.paragraphs
                ]
            shapes.append(entry)
        slides.append({"layout": slide.slide_layout.name, "shapes": shapes})
    return {
        "slide_count": len(prs.slides),
        "slides": slides,
    }


def main():
    if len(sys.argv) != 2:
        sys.exit("usage: validate_pptx.py <file.pptx>")
    print(json.dumps(summarize(sys.argv[1]), indent=2))


if __name__ == "__main__":
    main()
