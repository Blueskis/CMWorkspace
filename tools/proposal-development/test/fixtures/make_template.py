#!/usr/bin/env python3
"""Build a minimal but valid .potx fixture for testing, without needing a
confidential firm template. Five layouts (Title Slide, Title and Content,
Two Content, Section Header, Blank) plus a notes master, built with
python-pptx and then repackaged with a PresentationML template content type.

python-pptx cannot write a .potx content type directly, so this script builds
a normal .pptx first and then patches [Content_Types].xml and the extension.
"""
import shutil
import sys
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

HERE = Path(__file__).parent
OUT_PPTX = HERE / "_template_build.pptx"
OUT_POTX = HERE / "sample-template.potx"
OUT_PPTX_FINAL = HERE / "sample-template.pptx"


def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # python-pptx's default template already ships 11 layouts; we keep the
    # first five which match the names this fixture promises:
    # 0 Title Slide, 1 Title and Content, 2 Section Header, 3 Two Content,
    # 5 Title Only -> renamed conceptually as "Blank" usage via layout 6 Blank
    layouts = prs.slide_layouts
    wanted_idx = [0, 1, 2, 3, 6]  # Title Slide, Title+Content, Section Header, Two Content, Blank

    # Two example slides so example_slide_count == 2 (test case 2).
    slide = prs.slides.add_slide(layouts[0])
    slide.shapes.title.text = "Example Title Slide"
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Example subtitle"

    slide2 = prs.slides.add_slide(layouts[1])
    slide2.shapes.title.text = "Example Content Slide"

    # Touching notes_slide forces python-pptx to create a notes master,
    # so the fixture exercises the notesSupported: true path.
    slide.notes_slide.notes_text_frame.text = "Example speaker note."

    prs.save(str(OUT_PPTX))
    return wanted_idx


def patch_to_potx(wanted_idx):
    # Rezip, dropping any slideLayouts not in wanted_idx is unnecessary --
    # keep all layouts (real templates have many); the test only asserts the
    # five it cares about are readable with placeholders and geometry.
    shutil.copyfile(OUT_PPTX, OUT_PPTX_FINAL)

    with zipfile.ZipFile(OUT_PPTX, "r") as zin:
        names = zin.namelist()
        data = {n: zin.read(n) for n in names}

    ct = data["[Content_Types].xml"].decode("utf-8")
    ct = ct.replace(
        "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml",
        "application/vnd.openxmlformats-officedocument.presentationml.template.main+xml",
    )
    data["[Content_Types].xml"] = ct.encode("utf-8")

    with zipfile.ZipFile(OUT_POTX, "w", zipfile.ZIP_DEFLATED) as zout:
        for name, content in data.items():
            zout.writestr(name, content)

    OUT_PPTX.unlink()
    print(f"Wrote {OUT_POTX}")
    print(f"Wrote {OUT_PPTX_FINAL}")


def build_no_layouts_potx():
    """A degenerate .potx with a presentation.xml but no slideLayouts, for
    the zero-layout test case."""
    with zipfile.ZipFile(OUT_POTX, "r") as zin:
        data = {n: zin.read(n) for n in zin.namelist()}

    out = HERE / "empty-template.potx"
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zout:
        for name, content in data.items():
            if "slideLayouts" in name or "slideMasters" in name or "notesMaster" in name:
                continue
            zout.writestr(name, content)
    print(f"Wrote {out}")


def build_not_a_pptx():
    out = HERE / "not-a-template.potx"
    out.write_text("This is just a text file renamed to .potx\n")
    print(f"Wrote {out}")


def main():
    wanted_idx = build_pptx()
    patch_to_potx(wanted_idx)
    build_no_layouts_potx()
    build_not_a_pptx()
    print("Done.")


if __name__ == "__main__":
    sys.exit(main())
