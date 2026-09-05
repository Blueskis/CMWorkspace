#!/usr/bin/env python3
"""Assert no two top-level shapes on any slide overlap, and every shape sits within the
slide's bounds. Top-level only — a diagram's own internal boxes are laid out by
render_diagram's own non-overlapping geometry and are one grpSp here, not individually
checked. A small epsilon (EMU) tolerates shapes that intentionally share an edge (e.g. a
caption strip flush against its media rect) without flagging that as overlap.

This is the decisive guard for the body+picture composition fix: the reported bug was a
body block and a picture block landing on the IDENTICAL rectangle when a template has no
picture placeholder — this script is what would have caught that (100% overlap is still
overlap), run across every template in build-matrix.mjs's matrix.

    python3 test/check_overlap.py <built.pptx>
"""
import sys
from pptx import Presentation

EPS = 9144  # ~0.01in in EMU — tolerate shapes that intentionally share an edge


def rects_overlap(a, b, eps=EPS):
    ax1, ay1, ax2, ay2 = a
    bx1, by1, bx2, by2 = b
    return not (ax2 - eps <= bx1 or bx2 - eps <= ax1 or ay2 - eps <= by1 or by2 - eps <= ay1)


def main(path):
    prs = Presentation(path)
    sw, sh = prs.slide_width, prs.slide_height
    failures = []
    for si, slide in enumerate(prs.slides, 1):
        rects = []
        for shape in slide.shapes:
            if shape.left is None or shape.top is None:
                continue  # some placeholders with no explicit xfrm inherit from the layout; skip rather than false-positive
            l, t, w, h = shape.left, shape.top, shape.width, shape.height
            if l < -EPS or t < -EPS or l + w > sw + EPS or t + h > sh + EPS:
                failures.append(f"slide {si}: '{shape.name}' out of bounds ({l},{t},{w},{h}) vs slide {sw}x{sh}")
            rects.append((l, t, l + w, t + h, shape.name))
        for i in range(len(rects)):
            for j in range(i + 1, len(rects)):
                if rects_overlap(rects[i][:4], rects[j][:4]):
                    failures.append(f"slide {si}: '{rects[i][4]}' overlaps '{rects[j][4]}'")
    if failures:
        print("\n".join(failures))
        sys.exit(1)
    print("OK: no overlaps, all shapes in bounds")


if __name__ == "__main__":
    main(sys.argv[1])
